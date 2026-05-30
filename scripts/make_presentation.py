import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # High-end Color Palette (Cyber / Premium Dark mode)
    BG_COLOR = RGBColor(11, 10, 8)          # #0B0A08 (Space Black)
    SURFACE_COLOR = RGBColor(19, 18, 14)    # #13120E (Warm Charcoal Panel)
    TEXT_CREAM = RGBColor(244, 243, 239)    # #F4F3EF (Clean Warm White)
    TEXT_MUTED = RGBColor(156, 163, 175)    # #9CA3AF (Modern Muted Gray)
    ACCENT_GREEN = RGBColor(16, 185, 129)   # #10B981 (Cyber Emerald)
    ACCENT_COPPER = RGBColor(217, 119, 6)   # #D97706 (Warm Amber/Copper)
    ACCENT_MUTED = RGBColor(31, 41, 55)     # #1F2937 (Deep Slate)

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_layout_decorations(slide, category_text, title_text):
        # 1. Thin top accent line (Forest Green)
        top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = ACCENT_GREEN
        top_line.line.fill.background()

        # 2. Category tag (mono style)
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        ttf = tag_box.text_frame
        ttf.word_wrap = True
        ttf.margin_left = ttf.margin_top = ttf.margin_right = ttf.margin_bottom = 0
        tp = ttf.paragraphs[0]
        tp.text = f"AEOSPY  //  {category_text.upper()}"
        tp.font.name = "Courier New"
        tp.font.size = Pt(10)
        tp.font.bold = True
        tp.font.color.rgb = ACCENT_COPPER

        # 3. Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p_title = tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Georgia"
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = TEXT_CREAM

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide (Dark Theme)
    # ─────────────────────────────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, BG_COLOR)

    # Visual background structure (large accent blocks)
    accent_block = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = SURFACE_COLOR
    accent_block.line.fill.background()

    # Glowing radar ring visual mockup (stacked rings)
    for idx, radius in enumerate([1.8, 1.2, 0.6]):
        center_x = Inches(2.25) - Inches(radius/2)
        center_y = Inches(3.75) - Inches(radius/2)
        ring = slide1.shapes.add_shape(MSO_SHAPE.OVAL, center_x, center_y, Inches(radius), Inches(radius))
        ring.fill.background()
        ring.line.color.rgb = ACCENT_GREEN if idx == 0 else (ACCENT_COPPER if idx == 1 else TEXT_CREAM)
        ring.line.width = Pt(1.5 if idx == 2 else 0.8)

    # Scanner line mockup
    scan_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(2.0), Inches(0.05), Inches(1.75))
    scan_line.fill.solid()
    scan_line.fill.fore_color.rgb = ACCENT_GREEN
    scan_line.line.fill.background()

    # Title content
    title_box = slide1.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(7.333), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_tag = tf.paragraphs[0]
    p_tag.text = "PRODUCT PITCH  //  HACKATHON DECK"
    p_tag.font.name = "Courier New"
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_COPPER
    p_tag.space_after = Pt(16)

    p_main = tf.add_paragraph()
    p_main.text = "AEOspy"
    p_main.font.name = "Georgia"
    p_main.font.size = Pt(72)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_CREAM
    p_main.space_after = Pt(10)

    p_sub = tf.add_paragraph()
    p_sub.text = "The First Answer Engine Optimization Radar"
    p_sub.font.name = "Georgia"
    p_sub.font.size = Pt(24)
    p_sub.font.italic = True
    p_sub.font.color.rgb = ACCENT_GREEN
    p_sub.space_after = Pt(32)

    p_tech = tf.add_paragraph()
    p_tech.text = "POWERED BY BRIGHT DATA WEB DATA INTELLIGENCE"
    p_tech.font.name = "Courier New"
    p_tech.font.size = Pt(11)
    p_tech.font.color.rgb = TEXT_MUTED

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2: The Problem (The AEO Gap)
    # ─────────────────────────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, BG_COLOR)
    add_layout_decorations(slide2, "context", "The Problem: The AEO Gap")

    # Left Column: Stat Card
    stat_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    stat_bg.fill.solid()
    stat_bg.fill.fore_color.rgb = SURFACE_COLOR
    stat_bg.line.color.rgb = ACCENT_MUTED

    stat_box = slide2.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(4.6), Inches(4.0))
    stf = stat_box.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_top = stf.margin_right = stf.margin_bottom = 0

    sp1 = stf.paragraphs[0]
    sp1.text = "93%"
    sp1.font.name = "Georgia"
    sp1.font.size = Pt(84)
    sp1.font.bold = True
    sp1.font.color.rgb = ACCENT_COPPER
    sp1.space_after = Pt(4)

    sp2 = stf.add_paragraph()
    sp2.text = "of organic Google search leaders are completely omitted from ChatGPT and Perplexity recommendations."
    sp2.font.name = "Georgia"
    sp2.font.size = Pt(18)
    sp2.font.color.rgb = TEXT_CREAM
    sp2.space_after = Pt(16)

    sp3 = stf.add_paragraph()
    sp3.text = "Brands are investing heavily in SEO to rank #1 on Google, yet remaining completely invisible to the answer engines that consumers actually use to make decisions."
    sp3.font.name = "Arial"
    sp3.font.size = Pt(13)
    sp3.font.color.rgb = TEXT_MUTED

    # Right Column: Main Text & Cards
    right_box = slide2.shapes.add_textbox(Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.8))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = rtf.margin_top = rtf.margin_right = rtf.margin_bottom = 0

    rp1 = rtf.paragraphs[0]
    rp1.text = "CORE BARRIERS TO GTM REACH"
    rp1.font.name = "Courier New"
    rp1.font.size = Pt(11)
    rp1.font.bold = True
    rp1.font.color.rgb = ACCENT_GREEN
    rp1.space_after = Pt(16)

    challenges = [
        ("Black Box Citations", "LLMs pull facts and links from unmonitored indexes, offering marketing teams no dashboard or feedback loop to check visibility."),
        ("AI Hallucinations", "Generative engines fabricate brand constraints, features, or outdated price details, creating severe compliance and messaging risks."),
        ("Competitor Omissions", "In comparison queries, LLMs synthesize list comparisons, frequently prioritizing specific competitors due to structured schema layout alignment.")
    ]
    for idx, (title, desc) in enumerate(challenges):
        p_c = rtf.add_paragraph()
        p_c.text = f"0{idx+1}  {title}"
        p_c.font.name = "Georgia"
        p_c.font.size = Pt(16)
        p_c.font.bold = True
        p_c.font.color.rgb = TEXT_CREAM
        p_c.space_after = Pt(4)

        p_d = rtf.add_paragraph()
        p_d.text = desc
        p_d.font.name = "Arial"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(16)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3: The Solution (AEOspy)
    # ─────────────────────────────────────────────────────────────────────────
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, BG_COLOR)
    add_layout_decorations(slide3, "product", "The Solution: AEOspy Radar")

    # Column 1: Crawler
    c1_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8))
    c1_bg.fill.solid()
    c1_bg.fill.fore_color.rgb = SURFACE_COLOR
    c1_bg.line.color.rgb = ACCENT_MUTED

    c1_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.6), Inches(0.08))
    c1_line.fill.solid()
    c1_line.fill.fore_color.rgb = ACCENT_GREEN
    c1_line.line.fill.background()

    c1_box = slide3.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(3.0), Inches(4.0))
    c1_tf = c1_box.text_frame
    c1_tf.word_wrap = True
    c1_tf.margin_left = c1_tf.margin_top = c1_tf.margin_right = c1_tf.margin_bottom = 0
    c1_p1 = c1_tf.paragraphs[0]
    c1_p1.text = "01 / BRAND INDEX"
    c1_p1.font.name = "Courier New"
    c1_p1.font.size = Pt(11)
    c1_p1.font.bold = True
    c1_p1.font.color.rgb = ACCENT_COPPER
    c1_p1.space_after = Pt(12)
    
    c1_p2 = c1_tf.add_paragraph()
    c1_p2.text = "Multi-Agent Crawl"
    c1_p2.font.name = "Georgia"
    c1_p2.font.size = Pt(20)
    c1_p2.font.color.rgb = TEXT_CREAM
    c1_p2.space_after = Pt(12)

    c1_p3 = c1_tf.add_paragraph()
    c1_p3.text = "AEOspy extracts values, heading structures, and semantic entities directly from target homepages, creating the baseline profile."
    c1_p3.font.name = "Arial"
    c1_p3.font.size = Pt(13)
    c1_p3.font.color.rgb = TEXT_MUTED

    # Column 2: Audit
    c2_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    c2_bg.fill.solid()
    c2_bg.fill.fore_color.rgb = SURFACE_COLOR
    c2_bg.line.color.rgb = ACCENT_MUTED

    c2_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.6), Inches(0.08))
    c2_line.fill.solid()
    c2_line.fill.fore_color.rgb = ACCENT_GREEN
    c2_line.line.fill.background()

    c2_box = slide3.shapes.add_textbox(Inches(5.1), Inches(2.2), Inches(3.0), Inches(4.0))
    c2_tf = c2_box.text_frame
    c2_tf.word_wrap = True
    c2_tf.margin_left = c2_tf.margin_top = c2_tf.margin_right = c2_tf.margin_bottom = 0
    c2_p1 = c2_tf.paragraphs[0]
    c2_p1.text = "02 / ENGINE SCANS"
    c2_p1.font.name = "Courier New"
    c2_p1.font.size = Pt(11)
    c2_p1.font.bold = True
    c2_p1.font.color.rgb = ACCENT_COPPER
    c2_p1.space_after = Pt(12)

    c2_p2 = c2_tf.add_paragraph()
    c2_p2.text = "Citation Checks"
    c2_p2.font.name = "Georgia"
    c2_p2.font.size = Pt(20)
    c2_p2.font.color.rgb = TEXT_CREAM
    c2_p2.space_after = Pt(12)

    c2_p3 = c2_tf.add_paragraph()
    c2_p3.text = "Orchestrates parallel queries to 6 major AI engines using residential proxy nodes to capture actual user citation responses."
    c2_p3.font.name = "Arial"
    c2_p3.font.size = Pt(13)
    c2_p3.font.color.rgb = TEXT_MUTED

    # Column 3: Dashboard
    c3_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8))
    c3_bg.fill.solid()
    c3_bg.fill.fore_color.rgb = SURFACE_COLOR
    c3_bg.line.color.rgb = ACCENT_MUTED

    c3_line = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.6), Inches(0.08))
    c3_line.fill.solid()
    c3_line.fill.fore_color.rgb = ACCENT_GREEN
    c3_line.line.fill.background()

    c3_box = slide3.shapes.add_textbox(Inches(9.1), Inches(2.2), Inches(3.0), Inches(4.0))
    c3_tf = c3_box.text_frame
    c3_tf.word_wrap = True
    c3_tf.margin_left = c3_tf.margin_top = c3_tf.margin_right = c3_tf.margin_bottom = 0
    c3_p1 = c3_tf.paragraphs[0]
    c3_p1.text = "03 / INSIGHT HUD"
    c3_p1.font.name = "Courier New"
    c3_p1.font.size = Pt(11)
    c3_p1.font.bold = True
    c3_p1.font.color.rgb = ACCENT_COPPER
    c3_p1.space_after = Pt(12)

    c3_p2 = c3_tf.add_paragraph()
    c3_p2.text = "Gap Diagnostics"
    c3_p2.font.name = "Georgia"
    c3_p2.font.size = Pt(20)
    c3_p2.font.color.rgb = TEXT_CREAM
    c3_p2.space_after = Pt(12)

    c3_p3 = c3_tf.add_paragraph()
    c3_p3.text = "Generates citation indexes, flags hallucinations, suggests schema optimizations, and builds synthetic voice CMO summaries."
    c3_p3.font.name = "Arial"
    c3_p3.font.size = Pt(13)
    c3_p3.font.color.rgb = TEXT_MUTED

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4: Bright Data Integration
    # ─────────────────────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, BG_COLOR)
    add_layout_decorations(slide4, "data layer", "Bright Data Integration Flow")

    bd_features = [
        ("Web Unlocker", "CRAWLER ENGINE", "Bypasses proxy blocking to scrape brand homepages and formats raw content structures to clean markdown for downstream LLM synthesis."),
        ("SERP API", "ORGANIC INTELLIGENCE", "Pulls organic search rankings for target queries, allowing the system to run SEO vs AEO visibility gap comparisons."),
        ("Scraping Browser", "LLM SCRAPER", "Spawns automated Playwright browser instances over WebSocket CDP to query ChatGPT Search, Gemini, and Perplexity."),
        ("Dataset Scraping", "COMPETITOR SCRAPER", "Performs parallel batch scrapes of top-ranking competitor pages to extract strategic positioning and entity gaps.")
    ]

    for i, (name, tag, details) in enumerate(bd_features):
        left_pos = Inches(0.8 + i * 2.95)
        
        # Card Background
        bg = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.8))
        bg.fill.solid()
        bg.fill.fore_color.rgb = SURFACE_COLOR
        bg.line.color.rgb = ACCENT_MUTED

        # Left border highlight
        line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.8), Inches(0.06), Inches(4.8))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_GREEN
        line.line.fill.background()

        # Text Box
        box = slide4.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), Inches(2.4), Inches(4.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_tag = tf.paragraphs[0]
        p_tag.text = tag
        p_tag.font.name = "Courier New"
        p_tag.font.size = Pt(9)
        p_tag.font.bold = True
        p_tag.font.color.rgb = ACCENT_COPPER
        p_tag.space_after = Pt(10)
        
        p_name = tf.add_paragraph()
        p_name.text = name
        p_name.font.name = "Georgia"
        p_name.font.size = Pt(20)
        p_name.font.color.rgb = TEXT_CREAM
        p_name.space_after = Pt(12)

        p_desc = tf.add_paragraph()
        p_desc.text = details
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5: Multi-Model Sponsor Tracks
    # ─────────────────────────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, BG_COLOR)
    add_layout_decorations(slide5, "models", "Multi-Model Intelligence Pipeline")

    models = [
        ("Google Gemini (via AIML API)", "ENTITY PROCESSING & KEYWORDS", "Analyzes raw homepage markdown crawls to perform semantic entity extraction and construct list profiles of high-intent search keywords."),
        ("Anthropic Claude 3.5 Sonnet", "ADVANCED ANALYSIS SYNTHESIS", "Acts as the principal reasoning agent. Synthesizes crawler responses, ranking data, and citation structures into actionable recommendation streams via SSE."),
        ("Speechmatics Audio Engine", "SYNTHETIC CMO AUDIO BRIEFS", "Converts strategic insights into studio-quality executive briefings. Synthesizes a customizable CMO voice stream playable directly in the HUD.")
    ]

    for i, (provider, role, details) in enumerate(models):
        left_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8 + i * 1.65), Inches(11.73), Inches(1.4))
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = SURFACE_COLOR
        left_bg.line.color.rgb = ACCENT_MUTED

        # Left strip indicator
        strip = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8 + i * 1.65), Inches(0.08), Inches(1.4))
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT_GREEN if i == 1 else (ACCENT_COPPER if i == 0 else TEXT_CREAM)
        strip.line.fill.background()

        # Text Frame
        box = slide5.shapes.add_textbox(Inches(1.2), Inches(2.0 + i * 1.65), Inches(11.0), Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_role = tf.paragraphs[0]
        p_role.text = role
        p_role.font.name = "Courier New"
        p_role.font.size = Pt(9)
        p_role.font.bold = True
        p_role.font.color.rgb = ACCENT_COPPER
        p_role.space_after = Pt(4)

        p_title = tf.add_paragraph()
        p_title.text = provider
        p_title.font.name = "Georgia"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_CREAM
        
        run = p_title.add_run()
        run.text = f"  —  {details}"
        run.font.name = "Arial"
        run.font.size = Pt(13)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6: System Architecture
    # ─────────────────────────────────────────────────────────────────────────
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, BG_COLOR)
    add_layout_decorations(slide6, "architecture", "System Architecture")

    # Left: Text list
    left_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.8))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = ltf.margin_top = ltf.margin_right = ltf.margin_bottom = 0
    
    lp = ltf.paragraphs[0]
    lp.text = "Serverless Stream Infrastructure"
    lp.font.name = "Georgia"
    lp.font.size = Pt(22)
    lp.font.color.rgb = TEXT_CREAM
    lp.space_after = Pt(14)

    stack = [
        ("Next.js 15 App Router", "Server-side rendering with stream-compatible dynamic API routes."),
        ("Drizzle ORM & LibSQL/Turso", "Serverless-ready SQLite configuration that switches to cloud databases instantly via environment keys."),
        ("Server-Sent Events (SSE)", "Maintains persistent single-connection pipelines to push status streams directly to the dashboard HUD."),
        ("Framer Motion & Canvas LERP", "Highly responsive dark-mode layout utilizing lightweight math formulas to eliminate thread lag.")
    ]
    for title, desc in stack:
        p_s = ltf.add_paragraph()
        p_s.text = f"•  {title}: "
        p_s.font.name = "Georgia"
        p_s.font.size = Pt(13)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_COPPER
        
        run = p_s.add_run()
        run.text = desc
        run.font.name = "Arial"
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
        p_s.space_after = Pt(10)

    # Right: Pipeline Diagram
    diag_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(5.9), Inches(4.8))
    diag_bg.fill.solid()
    diag_bg.fill.fore_color.rgb = SURFACE_COLOR
    diag_bg.line.color.rgb = ACCENT_MUTED

    diag_box = slide6.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.1), Inches(4.2))
    dtf = diag_box.text_frame
    dtf.word_wrap = True
    dtf.margin_left = dtf.margin_top = dtf.margin_right = dtf.margin_bottom = 0

    p_d = dtf.paragraphs[0]
    p_d.text = "PIPELINE TOPOLOGY"
    p_d.font.name = "Courier New"
    p_d.font.size = Pt(11)
    p_d.font.bold = True
    p_d.font.color.rgb = ACCENT_GREEN
    p_d.space_after = Pt(16)

    steps = [
        ("Client Browser", "Initiates Audit Request via POST payload"),
        ("Next.js Server", "Logs database record & triggers AuditOrchestrator"),
        ("Orchestrator Swarm", "Deploys concurrent Web Scraper & SERP tasks"),
        ("Model Pipeline", "Gemini processes keywords; Claude builds final analysis"),
        ("SSE Controller", "Streams live progress updates back to Client Browser")
    ]
    for idx, (step, desc) in enumerate(steps):
        p_step = dtf.add_paragraph()
        p_step.text = f"0{idx+1}  {step}"
        p_step.font.name = "Georgia"
        p_step.font.size = Pt(14)
        p_step.font.bold = True
        p_step.font.color.rgb = TEXT_CREAM
        p_step.space_after = Pt(3)
        
        p_step_desc = dtf.add_paragraph()
        p_step_desc.text = desc
        p_step_desc.font.name = "Arial"
        p_step_desc.font.size = Pt(11)
        p_step_desc.font.color.rgb = TEXT_MUTED
        p_step_desc.space_after = Pt(10)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Summary & Pitch
    # ─────────────────────────────────────────────────────────────────────────
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, SURFACE_COLOR)

    # Glowing center frame
    pitch_bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1))
    pitch_bg.fill.solid()
    pitch_bg.fill.fore_color.rgb = BG_COLOR
    pitch_bg.line.color.rgb = ACCENT_GREEN

    pitch_box = slide7.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9.73), Inches(4.5))
    ptf = pitch_box.text_frame
    ptf.word_wrap = True
    ptf.margin_left = ptf.margin_top = ptf.margin_right = ptf.margin_bottom = 0

    pp1 = ptf.paragraphs[0]
    pp1.text = "AEOspy"
    pp1.font.name = "Georgia"
    pp1.font.size = Pt(48)
    pp1.font.bold = True
    pp1.font.color.rgb = TEXT_CREAM
    pp1.space_after = Pt(8)
    pp1.alignment = PP_ALIGN.CENTER

    pp2 = ptf.add_paragraph()
    pp2.text = "The Future of Web Visibility is Answer Engine Optimization"
    pp2.font.name = "Georgia"
    pp2.font.size = Pt(24)
    pp2.font.italic = True
    pp2.font.color.rgb = ACCENT_COPPER
    pp2.space_after = Pt(24)
    pp2.alignment = PP_ALIGN.CENTER

    pp3 = ptf.add_paragraph()
    pp3.text = "By unifying Bright Data's Web Scraping proxies with Multi-Model Intelligence, AEOspy provides brands with the tools to transition from classic SEO to the era of Generative AI Search."
    pp3.font.name = "Arial"
    pp3.font.size = Pt(16)
    pp3.font.color.rgb = TEXT_MUTED
    pp3.space_after = Pt(36)
    pp3.alignment = PP_ALIGN.CENTER

    pp4 = ptf.add_paragraph()
    pp4.text = "https://github.com/Shikhyy/AEOspy"
    pp4.font.name = "Courier New"
    pp4.font.size = Pt(14)
    pp4.font.bold = True
    pp4.font.color.rgb = ACCENT_GREEN
    pp4.alignment = PP_ALIGN.CENTER

    prs.save("/Users/shikhar/AEOspy/AEOspy_Pitch_Deck.pptx")
    print("Presentation upgraded successfully at /Users/shikhar/AEOspy/AEOspy_Pitch_Deck.pptx!")

if __name__ == "__main__":
    create_presentation()
